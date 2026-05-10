"""
artifact_export.py — Artifact を PDF / Excel / PPTX に書き出す。

各 view 型 × 各形式の組み合わせで、テンプレ（minimal / corporate / branded）を選べる。
最初は minimal だけ。テンプレ追加は assets/templates/<format>/<template>.json。

責務:
  - artifact + format + template → ファイル
  - 出力先: <repo>/data/exports/<artifact_id>/<timestamp>.<ext>
  - URL は /api/artifacts/{id}/exports/{filename} で配信

依存（pip）:
  openpyxl / python-pptx / reportlab
"""

from __future__ import annotations

import io
import json
from datetime import datetime
from pathlib import Path
from typing import Optional

EXPORT_DIR = Path(__file__).resolve().parents[2] / "data" / "exports"
EXPORT_DIR.mkdir(parents=True, exist_ok=True)


def _ts() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


def _output_path(artifact_id: str, ext: str) -> Path:
    d = EXPORT_DIR / artifact_id
    d.mkdir(parents=True, exist_ok=True)
    return d / f"{_ts()}.{ext}"


# ──────────────────────────────────────────
# Excel エクスポート
# ──────────────────────────────────────────

def export_to_excel(artifact: dict, template: str = "minimal") -> Path:
    """artifact → .xlsx
    対応 view: list / table / kanban / kpi-card / compare / form / matrix
    その他は JSON ダンプを 1 セルに置く（fallback）"""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment

    wb = Workbook()
    ws = wb.active
    ws.title = artifact["type"]

    title_font = Font(bold=True, size=14)
    head_font = Font(bold=True)
    head_fill = PatternFill("solid", fgColor="E5E7EB")

    ws["A1"] = artifact.get("title") or artifact["type"]
    ws["A1"].font = title_font

    t = artifact["type"]
    data = artifact.get("data") or {}

    if t == "list":
        ws["A3"] = "Done"; ws["B3"] = "Item"
        ws["A3"].font = head_font; ws["B3"].font = head_font
        for i, it in enumerate(data.get("items") or []):
            ws.cell(row=4 + i, column=1).value = "Yes" if it.get("done") else ""
            ws.cell(row=4 + i, column=2).value = it.get("text", "")
        ws.column_dimensions["B"].width = 60

    elif t == "table":
        cols = data.get("columns") or []
        for j, c in enumerate(cols):
            cell = ws.cell(row=3, column=j + 1, value=c)
            cell.font = head_font; cell.fill = head_fill
        for i, row in enumerate(data.get("rows") or []):
            for j, cell in enumerate(row):
                ws.cell(row=4 + i, column=j + 1, value=cell)
        for j in range(len(cols)):
            ws.column_dimensions[chr(65 + j)].width = 20

    elif t == "kanban":
        cols = data.get("columns") or []
        for j, col in enumerate(cols):
            cell = ws.cell(row=3, column=j + 1, value=col.get("title", col.get("id")))
            cell.font = head_font; cell.fill = head_fill
        max_cards = max((len(c.get("cards") or []) for c in cols), default=0)
        for j, col in enumerate(cols):
            for i, card in enumerate(col.get("cards") or []):
                ws.cell(row=4 + i, column=j + 1, value=card.get("text", ""))
        for j in range(len(cols)):
            ws.column_dimensions[chr(65 + j)].width = 30

    elif t == "kpi-card":
        ws["A3"] = "指標"; ws["B3"] = "値"; ws["C3"] = "単位"; ws["D3"] = "前期差"
        for c in ("A3", "B3", "C3", "D3"):
            ws[c].font = head_font; ws[c].fill = head_fill
        metrics = data.get("metrics") or []
        for i, m in enumerate(metrics):
            ws.cell(row=4 + i, column=1, value=m.get("label", ""))
            ws.cell(row=4 + i, column=2, value=m.get("value"))
            ws.cell(row=4 + i, column=3, value=m.get("unit", ""))
            ws.cell(row=4 + i, column=4, value=m.get("delta"))

    elif t == "compare":
        items = data.get("items") or []
        criteria = data.get("criteria") or (
            list((items[0].get("values") or {}).keys()) if items else []
        )
        ws["A3"] = "項目"; ws["A3"].font = head_font; ws["A3"].fill = head_fill
        for j, it in enumerate(items):
            cell = ws.cell(row=3, column=2 + j, value=it.get("name", f"案{j+1}"))
            cell.font = head_font; cell.fill = head_fill
        for i, c in enumerate(criteria):
            ws.cell(row=4 + i, column=1, value=c).font = head_font
            for j, it in enumerate(items):
                v = (it.get("values") or {}).get(c)
                ws.cell(row=4 + i, column=2 + j, value=str(v) if v is not None else "")

    elif t == "form":
        fields = data.get("fields") or []
        values = data.get("values") or {}
        ws["A3"] = "項目"; ws["B3"] = "回答"
        ws["A3"].font = head_font; ws["B3"].font = head_font
        for i, f in enumerate(fields):
            ws.cell(row=4 + i, column=1, value=f.get("label", f.get("id", "")))
            ws.cell(row=4 + i, column=2, value=str(values.get(f.get("id"), "")))
        ws.column_dimensions["B"].width = 50

    elif t == "matrix":
        labels = data.get("labels") or {
            "q1": "重要・緊急", "q2": "重要・非緊急",
            "q3": "緊急・非重要", "q4": "非緊急・非重要",
        }
        ws.cell(row=3, column=1, value=labels.get("q1", "")).font = head_font
        ws.cell(row=3, column=2, value=labels.get("q2", "")).font = head_font
        ws.cell(row=10, column=1, value=labels.get("q3", "")).font = head_font
        ws.cell(row=10, column=2, value=labels.get("q4", "")).font = head_font
        items_by_q = {"q1": [], "q2": [], "q3": [], "q4": []}
        for it in data.get("items") or []:
            items_by_q.setdefault(it.get("quadrant", "q1"), []).append(it.get("text", ""))
        for q, col_row in [("q1", (4, 1)), ("q2", (4, 2)), ("q3", (11, 1)), ("q4", (11, 2))]:
            for i, txt in enumerate(items_by_q[q]):
                ws.cell(row=col_row[0] + i, column=col_row[1], value=txt)

    else:
        ws["A3"] = "data"
        ws["A4"] = json.dumps(data, ensure_ascii=False, indent=2)
        ws["A4"].alignment = Alignment(wrap_text=True, vertical="top")
        ws.column_dimensions["A"].width = 80

    out = _output_path(artifact["id"], "xlsx")
    wb.save(out)
    return out


# ──────────────────────────────────────────
# PowerPoint エクスポート
# ──────────────────────────────────────────

def export_to_pptx(artifact: dict, template: str = "minimal") -> Path:
    """artifact → .pptx
    対応 view: slide / list / table / kanban / kpi-card / compare
    その他は 1 枚にタイトル + JSON 縮小表示"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    def add_title_slide(title: str, sub: str = ""):
        s = prs.slides.add_slide(blank)
        tx = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5)).text_frame
        tx.text = title
        tx.paragraphs[0].alignment = PP_ALIGN.CENTER
        tx.paragraphs[0].runs[0].font.size = Pt(44)
        tx.paragraphs[0].runs[0].font.bold = True
        if sub:
            sx = s.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(1)).text_frame
            sx.text = sub
            sx.paragraphs[0].alignment = PP_ALIGN.CENTER
            sx.paragraphs[0].runs[0].font.size = Pt(20)
        return s

    def add_content_slide(title: str, lines: list[str]):
        s = prs.slides.add_slide(blank)
        t = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8)).text_frame
        t.text = title
        t.paragraphs[0].runs[0].font.size = Pt(28)
        t.paragraphs[0].runs[0].font.bold = True
        b = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5)).text_frame
        b.word_wrap = True
        for i, ln in enumerate(lines):
            p = b.paragraphs[0] if i == 0 else b.add_paragraph()
            p.text = ln
            p.runs[0].font.size = Pt(18)
        return s

    title = artifact.get("title") or artifact["type"]
    t = artifact["type"]
    data = artifact.get("data") or {}

    add_title_slide(title, f"({t})")

    if t == "slide":
        for sl in data.get("slides") or []:
            add_content_slide(sl.get("title") or "", [sl.get("body", "")])

    elif t == "list":
        items = data.get("items") or []
        lines = [
            ("[x] " if it.get("done") else "[ ] ") + (it.get("text") or "")
            for it in items
        ]
        add_content_slide("項目", lines)

    elif t == "table":
        cols = data.get("columns") or []
        rows = data.get("rows") or []
        s = prs.slides.add_slide(blank)
        ttl = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7)).text_frame
        ttl.text = title; ttl.paragraphs[0].runs[0].font.size = Pt(24); ttl.paragraphs[0].runs[0].font.bold = True
        if cols and rows:
            tbl = s.shapes.add_table(
                rows=len(rows) + 1, cols=len(cols),
                left=Inches(0.5), top=Inches(1.2), width=Inches(12), height=Inches(5),
            ).table
            for j, c in enumerate(cols):
                cell = tbl.cell(0, j); cell.text = c
                for r in cell.text_frame.paragraphs[0].runs:
                    r.font.bold = True
            for i, row in enumerate(rows):
                for j, val in enumerate(row[: len(cols)]):
                    tbl.cell(i + 1, j).text = str(val)

    elif t == "kanban":
        cols = data.get("columns") or []
        s = prs.slides.add_slide(blank)
        ttl = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7)).text_frame
        ttl.text = title; ttl.paragraphs[0].runs[0].font.size = Pt(24); ttl.paragraphs[0].runs[0].font.bold = True
        col_w = 12 / max(1, len(cols))
        for j, col in enumerate(cols):
            x = Inches(0.5 + j * col_w)
            box = s.shapes.add_textbox(x, Inches(1.3), Inches(col_w - 0.2), Inches(5.5)).text_frame
            box.word_wrap = True
            box.text = col.get("title", "")
            box.paragraphs[0].runs[0].font.bold = True; box.paragraphs[0].runs[0].font.size = Pt(16)
            for cd in col.get("cards") or []:
                p = box.add_paragraph(); p.text = "・" + (cd.get("text") or "")
                p.runs[0].font.size = Pt(12)

    elif t == "kpi-card":
        metrics = data.get("metrics") or []
        s = prs.slides.add_slide(blank)
        ttl = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7)).text_frame
        ttl.text = title; ttl.paragraphs[0].runs[0].font.size = Pt(24); ttl.paragraphs[0].runs[0].font.bold = True
        cards = metrics[:6]
        for i, m in enumerate(cards):
            row, col = divmod(i, 3)
            x = Inches(0.5 + col * 4.2)
            y = Inches(1.3 + row * 2.8)
            tx = s.shapes.add_textbox(x, y, Inches(4), Inches(2.5)).text_frame
            tx.text = m.get("label", "")
            tx.paragraphs[0].runs[0].font.size = Pt(14)
            tx.paragraphs[0].runs[0].font.color.rgb = None
            v = tx.add_paragraph()
            v.text = f"{m.get('value','')}{m.get('unit','')}"
            v.runs[0].font.size = Pt(36); v.runs[0].font.bold = True
            if m.get("delta") is not None:
                d = tx.add_paragraph()
                d.text = ("+" if (m.get("delta") or 0) > 0 else "-") + str(abs(m.get("delta", 0)))
                d.runs[0].font.size = Pt(12)

    elif t == "compare":
        items = data.get("items") or []
        criteria = data.get("criteria") or (
            list((items[0].get("values") or {}).keys()) if items else []
        )
        s = prs.slides.add_slide(blank)
        ttl = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7)).text_frame
        ttl.text = title; ttl.paragraphs[0].runs[0].font.size = Pt(24); ttl.paragraphs[0].runs[0].font.bold = True
        if items and criteria:
            tbl = s.shapes.add_table(
                rows=len(criteria) + 1, cols=len(items) + 1,
                left=Inches(0.5), top=Inches(1.2), width=Inches(12), height=Inches(5),
            ).table
            tbl.cell(0, 0).text = "項目"
            for j, it in enumerate(items):
                tbl.cell(0, j + 1).text = it.get("name", "")
            for i, c in enumerate(criteria):
                tbl.cell(i + 1, 0).text = c
                for j, it in enumerate(items):
                    v = (it.get("values") or {}).get(c)
                    tbl.cell(i + 1, j + 1).text = str(v) if v is not None else ""

    else:
        add_content_slide("data", [json.dumps(data, ensure_ascii=False)[:1500]])

    out = _output_path(artifact["id"], "pptx")
    prs.save(out)
    return out


# ──────────────────────────────────────────
# PDF エクスポート（reportlab）
# ──────────────────────────────────────────

def export_to_pdf(artifact: dict, template: str = "minimal") -> Path:
    """artifact → .pdf（A4 縦・基本日本語フォントで描画）"""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, ListFlowable, ListItem,
    )

    # 日本語埋め込みフォント（macOS 標準 + reportlab CID）
    try:
        pdfmetrics.registerFont(UnicodeCIDFont("HeiseiKakuGo-W5"))
        ja = "HeiseiKakuGo-W5"
    except Exception:
        ja = "Helvetica"

    out = _output_path(artifact["id"], "pdf")
    doc = SimpleDocTemplate(str(out), pagesize=A4, leftMargin=20*mm, rightMargin=20*mm)
    styles = getSampleStyleSheet()
    h1 = ParagraphStyle("H1", parent=styles["Heading1"], fontName=ja, fontSize=18, leading=22)
    h2 = ParagraphStyle("H2", parent=styles["Heading2"], fontName=ja, fontSize=13, leading=17)
    body = ParagraphStyle("B", parent=styles["BodyText"], fontName=ja, fontSize=10, leading=14)

    story = []
    story.append(Paragraph(artifact.get("title") or artifact["type"], h1))
    story.append(Paragraph(f"({artifact['type']})", body))
    story.append(Spacer(1, 8))

    t = artifact["type"]
    data = artifact.get("data") or {}

    if t == "list":
        items = [
            ("[x] " if it.get("done") else "[ ] ") + (it.get("text") or "")
            for it in (data.get("items") or [])
        ]
        story.append(ListFlowable([ListItem(Paragraph(x, body)) for x in items], bulletType="bullet"))

    elif t == "table" and data.get("columns") and data.get("rows"):
        rows = [data["columns"], *data["rows"]]
        tbl = Table(rows, repeatRows=1)
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E5E7EB")),
            ("FONTNAME", (0, 0), (-1, -1), ja),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ]))
        story.append(tbl)

    elif t == "kpi-card":
        metrics = data.get("metrics") or []
        rows = [["指標", "値", "単位", "前期差"]]
        for m in metrics:
            rows.append([
                m.get("label", ""), str(m.get("value", "")),
                m.get("unit", ""), str(m.get("delta", "")),
            ])
        tbl = Table(rows, repeatRows=1)
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E5E7EB")),
            ("FONTNAME", (0, 0), (-1, -1), ja),
            ("FONTSIZE", (0, 0), (-1, -1), 10),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
        ]))
        story.append(tbl)

    elif t == "markdown":
        text = (data.get("text") or "")[:8000]
        for line in text.split("\n"):
            if line.startswith("# "):
                story.append(Paragraph(line[2:], h1))
            elif line.startswith("## "):
                story.append(Paragraph(line[3:], h2))
            elif line.strip():
                story.append(Paragraph(line.replace("&", "&amp;").replace("<", "&lt;"), body))
            else:
                story.append(Spacer(1, 6))

    elif t == "compare":
        items = data.get("items") or []
        criteria = data.get("criteria") or (
            list((items[0].get("values") or {}).keys()) if items else []
        )
        rows = [["項目"] + [it.get("name", "") for it in items]]
        for c in criteria:
            rows.append([c] + [str((it.get("values") or {}).get(c, "")) for it in items])
        tbl = Table(rows, repeatRows=1)
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E5E7EB")),
            ("FONTNAME", (0, 0), (-1, -1), ja),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
        ]))
        story.append(tbl)

    elif t == "kanban":
        for col in data.get("columns") or []:
            story.append(Paragraph(col.get("title", ""), h2))
            cards = col.get("cards") or []
            for cd in cards:
                story.append(Paragraph("・" + (cd.get("text") or ""), body))
            story.append(Spacer(1, 8))

    elif t == "form":
        fields = data.get("fields") or []
        values = data.get("values") or {}
        rows = [["項目", "回答"]]
        for f in fields:
            rows.append([
                f.get("label", f.get("id", "")),
                str(values.get(f.get("id"), "")),
            ])
        tbl = Table(rows, repeatRows=1)
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E5E7EB")),
            ("FONTNAME", (0, 0), (-1, -1), ja),
            ("FONTSIZE", (0, 0), (-1, -1), 10),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
        ]))
        story.append(tbl)

    else:
        story.append(Paragraph(f"<pre>{json.dumps(data, ensure_ascii=False, indent=2)[:4000]}</pre>", body))

    doc.build(story)
    return out


# ──────────────────────────────────────────
# 統合ディスパッチャ
# ──────────────────────────────────────────

EXPORTERS = {
    "pdf":  export_to_pdf,
    "xlsx": export_to_excel,
    "excel": export_to_excel,
    "pptx": export_to_pptx,
    "ppt":  export_to_pptx,
}


def export_artifact(artifact: dict, format: str, template: str = "minimal") -> Path:
    fmt = format.lower().lstrip(".")
    fn = EXPORTERS.get(fmt)
    if not fn:
        raise ValueError(f"unsupported format: {format}")
    return fn(artifact, template=template)
